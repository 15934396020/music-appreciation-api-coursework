"""Generate a professional PPTX presentation for the Music Appreciation and Discovery API coursework.

This script creates a 12-slide PowerPoint presentation covering all required
deliverable sections: project overview, architecture, implementation, testing,
version control, API documentation, technical report highlights, and Q&A.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colour palette
DARK_BLUE = RGBColor(0x1A, 0x47, 0x7A)
MEDIUM_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE = RGBColor(0xD4, 0xE6, 0xF1)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_GREY = RGBColor(0xEC, 0xF0, 0xF1)
BLACK = RGBColor(0x00, 0x00, 0x00)


def add_background(slide, color=WHITE):
    """Set slide background colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a coloured title bar at the top of the slide."""
    # Title bar shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = LIGHT_BLUE
        p2.alignment = PP_ALIGN.LEFT


def add_bullet_content(slide, items, left=0.6, top=1.5, width=12, font_size=18, spacing=Pt(8)):
    """Add bullet-point content to a slide."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            # (text, is_bold, color)
            p.text = item[0]
            p.font.bold = item[1] if len(item) > 1 else False
            p.font.color.rgb = item[2] if len(item) > 2 else DARK_GREY
            p.font.size = Pt(item[3]) if len(item) > 3 else Pt(font_size)
        else:
            p.text = f"  {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK_GREY

        p.space_after = spacing
        p.alignment = PP_ALIGN.LEFT


def add_table(slide, headers, rows, left=0.6, top=None, width=12, col_widths=None):
    """Add a formatted table to the slide."""
    if top is None:
        top = 2.0
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(0.4 * num_rows),
    )
    table = table_shape.table

    # Set column widths if provided
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    # Data rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = DARK_GREY
                paragraph.alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GREY

    return table_shape


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # -----------------------------------------------------------------------
    # Slide 1: Title Slide
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, DARK_BLUE)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Music Appreciation and Discovery API"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "XJCO3011 — Web Services and Web Data"
    p2.font.size = Pt(22)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Coursework 1: Individual Web Services API Development Project"
    p3.font.size = Pt(18)
    p3.font.color.rgb = LIGHT_BLUE
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(10)

    # Bottom info
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p4 = tf2.paragraphs[0]
    p4.text = "FastAPI  |  SQLAlchemy  |  SQLite  |  Pydantic  |  pytest"
    p4.font.size = Pt(16)
    p4.font.color.rgb = ACCENT_ORANGE
    p4.alignment = PP_ALIGN.CENTER

    p5 = tf2.add_paragraph()
    p5.text = "25 Endpoints  •  55 Tests  •  API Key Authentication  •  5 Analytics Endpoints"
    p5.font.size = Pt(14)
    p5.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    p5.alignment = PP_ALIGN.CENTER
    p5.space_before = Pt(10)

    # -----------------------------------------------------------------------
    # Slide 2: Problem Statement & Motivation
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Problem Statement & Motivation")

    add_bullet_content(slide, [
        ("The Challenge", True, MEDIUM_BLUE, 22),
        "Music enthusiasts lack a structured way to catalogue listening experiences,",
        "group tracks by mood or genre, and share meaningful reviews.",
        "",
        ("Our Solution", True, MEDIUM_BLUE, 22),
        "A RESTful API that provides structured music metadata browsing,",
        "full CRUD review management, user-defined tagging, collection building,",
        "and analytical insights — all backed by a relational database.",
        "",
        ("Design Philosophy", True, MEDIUM_BLUE, 22),
        "Deliberate scope control: focus on stability, testability, and explainability",
        "rather than maximum technical complexity.",
    ], top=1.5)

    # -----------------------------------------------------------------------
    # Slide 3: Scope Control
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Project Scope")

    add_table(slide,
        ["In Scope", "Out of Scope"],
        [
            ["29 tracks across 8 genres (curated seed data)", "Audio playback / streaming"],
            ["Full CRUD for reviews (Create, Read, Update, Delete)", "Recommendation engine / ML"],
            ["User tags and collection management", "External API integration (Spotify)"],
            ["5 analytics endpoints with SQL aggregation", "Cloud deployment (future work)"],
            ["API key authentication for write operations", "Per-user OAuth2/JWT (future work)"],
            ["55 automated tests across 9 test classes", "Frontend UI"],
            ["Structured JSON error responses", "Rate limiting (future work)"],
        ],
        top=1.6,
        col_widths=[6, 6],
    )

    # -----------------------------------------------------------------------
    # Slide 4: Architecture & Technology Stack
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Architecture & Technology Stack")

    add_table(slide,
        ["Layer", "Technology", "Justification"],
        [
            ["Client", "HTTP / Swagger UI", "Interactive testing and documentation"],
            ["Authentication", "API Key (X-API-Key header)", "Lightweight access control for write operations"],
            ["Routing", "FastAPI 0.115.12", "High performance, auto OpenAPI docs, type hints"],
            ["Validation", "Pydantic 2.11.3", "Strict payload validation with clear error messages"],
            ["ORM", "SQLAlchemy 2.0.40", "Powerful querying for analytics without raw SQL"],
            ["Database", "SQLite", "Zero-config setup, easy for examiners to run"],
            ["Testing", "pytest + httpx", "Session-scoped fixtures, TestClient integration"],
        ],
        top=1.6,
        col_widths=[2.5, 4, 5.5],
    )

    # -----------------------------------------------------------------------
    # Slide 5: Data Model
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Database Schema: 6 Entities")

    add_table(slide,
        ["Entity", "Key Fields", "Relationships"],
        [
            ["Genre", "id, name, description", "One-to-many with Track"],
            ["Track", "id, title, artist, album, mood, avg_rating", "Belongs to Genre; has Reviews, Tags, CollectionItems"],
            ["Review", "id, track_id, reviewer, rating (1-5), comment", "Belongs to Track; triggers rating recalculation"],
            ["UserTag", "id, track_id, tag_name, created_by", "Belongs to Track; unique(track_id, tag_name)"],
            ["Collection", "id, name, description, created_by", "Has many CollectionItems"],
            ["CollectionItem", "id, collection_id, track_id, note", "Links Collection to Track; unique(collection_id, track_id)"],
        ],
        top=1.6,
        col_widths=[2.5, 5, 4.5],
    )

    # -----------------------------------------------------------------------
    # Slide 6: API Endpoints Overview
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "API Endpoints: 25 Routes Across 7 Groups")

    add_table(slide,
        ["Group", "Count", "Key Endpoints", "Auth Required"],
        [
            ["General", "2", "GET /, GET /health", "No"],
            ["Genres", "2", "GET /genres, GET /genres/{id}", "No"],
            ["Tracks", "2", "GET /tracks (filters + pagination), GET /tracks/{id}", "No"],
            ["Reviews", "5", "POST, GET (list), GET (by ID), PUT, DELETE", "POST/PUT/DELETE"],
            ["Tags", "3", "POST /tags, GET /tags, DELETE /tags/{id}", "POST/DELETE"],
            ["Collections", "6", "POST, GET (list), GET (detail), POST (add item), DELETE (item), DELETE (collection)", "POST/DELETE"],
            ["Analytics", "5", "top-rated-tracks, genre-summary, top-tags, mood-distribution, review-activity", "No"],
        ],
        top=1.6,
        col_widths=[2, 1.2, 7, 1.8],
    )

    # -----------------------------------------------------------------------
    # Slide 7: Authentication & Error Handling
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Authentication & Error Handling")

    add_bullet_content(slide, [
        ("API Key Authentication", True, MEDIUM_BLUE, 22),
        "Write operations (POST, PUT, DELETE) require X-API-Key header",
        "Read operations (GET) remain publicly accessible",
        "401 Unauthorized for missing key, 403 Forbidden for invalid key",
        "Demo key: music-api-demo-key-2026",
        "",
        ("Structured Error Responses", True, MEDIUM_BLUE, 22),
        'All errors return consistent JSON: {"error": "...", "message": "...", "details": [...]}',
        "Custom exception handlers for HTTP errors and validation errors",
        "Machine-readable error codes (not_found, validation_error, conflict, etc.)",
        "",
        ("Status Codes Used", True, MEDIUM_BLUE, 22),
        "200 OK, 201 Created, 204 No Content, 401 Unauthorized, 403 Forbidden, 404 Not Found, 409 Conflict, 422 Validation Error",
    ], top=1.5)

    # -----------------------------------------------------------------------
    # Slide 8: Review CRUD Demo
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Review CRUD: Full Lifecycle Demo")

    add_bullet_content(slide, [
        ("CREATE — POST /reviews (with X-API-Key)", True, ACCENT_GREEN, 20),
        '  Request: {"track_id": 1, "reviewer_name": "Alice", "rating": 5, "comment": "A masterpiece"}',
        "  Response: 201 Created with full review object including auto-generated timestamps",
        "",
        ("READ — GET /reviews?track_id=1&min_rating=4", True, MEDIUM_BLUE, 20),
        "  Supports filtering by track_id, reviewer_name, min_rating with pagination",
        "",
        ("UPDATE — PUT /reviews/{id} (with X-API-Key)", True, ACCENT_ORANGE, 20),
        '  Request: {"rating": 4, "comment": "Updated opinion"} — partial update supported',
        "  Automatically recalculates track average_rating via _refresh_track_rating()",
        "",
        ("DELETE — DELETE /reviews/{id} (with X-API-Key)", True, ACCENT_RED, 20),
        "  Returns 204 No Content; also triggers average rating recalculation",
    ], top=1.5, font_size=16)

    # -----------------------------------------------------------------------
    # Slide 9: Analytics Endpoints
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Analytics: 5 Aggregation Endpoints")

    add_table(slide,
        ["Endpoint", "SQLAlchemy Features", "Returns"],
        [
            ["/analytics/top-rated-tracks", "func.avg, func.count, group_by, having, outerjoin", "Ranked tracks with review counts"],
            ["/analytics/genre-summary", "func.count, func.avg, outerjoin, group_by", "Track count and avg rating per genre"],
            ["/analytics/top-tags", "func.count, group_by, order_by", "Most frequently used tags"],
            ["/analytics/mood-distribution", "func.count, group_by, filter", "Track count per mood category"],
            ["/analytics/review-activity", "func.count, func.avg, group_by", "Per-reviewer stats and totals"],
        ],
        top=1.6,
        col_widths=[3.5, 5, 3.5],
    )

    add_bullet_content(slide, [
        "",
        ("Key Insight:", True, ACCENT_ORANGE, 18),
        "These endpoints demonstrate advanced SQL aggregation through SQLAlchemy ORM,",
        "elevating the project beyond basic CRUD to provide real analytical value.",
    ], top=4.8, font_size=16)

    # -----------------------------------------------------------------------
    # Slide 10: Testing & Reliability
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Testing: 55 Automated Tests, 9 Test Classes")

    add_table(slide,
        ["Test Class", "Tests", "Coverage Area"],
        [
            ["TestGeneralEndpoints", "2", "Root and health endpoints"],
            ["TestGenreEndpoints", "3", "Genre listing and lookup"],
            ["TestTrackEndpoints", "8", "Track browsing, filtering, pagination"],
            ["TestReviewCRUD", "9", "Full review lifecycle + error cases"],
            ["TestTagEndpoints", "7", "Tag CRUD + duplicate/invalid handling"],
            ["TestCollectionEndpoints", "9", "Collection management + item operations"],
            ["TestAnalyticsEndpoints", "5", "All 5 analytics endpoints"],
            ["TestAuthentication", "6", "API key validation (401, 403, public access)"],
            ["TestValidation", "6", "Input validation edge cases + structured errors"],
        ],
        top=1.6,
        col_widths=[3.5, 1.2, 7.3],
    )

    add_bullet_content(slide, [
        ("All 55 tests pass in under 1 second using pytest with session-scoped TestClient fixture.", False, ACCENT_GREEN, 16),
    ], top=6.2, font_size=16)

    # -----------------------------------------------------------------------
    # Slide 11: Version Control & Deliverables
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Version Control & Deliverables")

    add_bullet_content(slide, [
        ("GitHub Repository", True, MEDIUM_BLUE, 22),
        "Public repository: github.com/15934396020/music-appreciation-api-coursework",
        "Consistent commit history showing incremental development",
        "Clean branch structure (main branch)",
        "",
        ("Deliverables Checklist", True, MEDIUM_BLUE, 22),
        "Technical Report (PDF) — 5 pages covering design, implementation, testing, GenAI declaration",
        "API Documentation (PDF) — Complete endpoint reference with examples",
        "Presentation Slides (PPTX) — This presentation",
        "GenAI Conversation Log — Detailed AI usage documentation",
        "README.md — Setup instructions, endpoint overview, project structure",
        "",
        ("Multi-Account Handover Protocol", True, MEDIUM_BLUE, 22),
        "Structured handover documents enabling seamless AI-assisted development across sessions",
    ], top=1.5)

    # -----------------------------------------------------------------------
    # Slide 12: Conclusion & Future Work
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Conclusion & Future Work")

    # Three columns
    # Strengths
    add_bullet_content(slide, [
        ("Strengths", True, ACCENT_GREEN, 22),
        "Clean, modular FastAPI architecture",
        "25 endpoints with full CRUD",
        "API key authentication",
        "Structured error handling",
        "55 passing automated tests",
        "5 analytics endpoints with SQL aggregation",
        "Comprehensive documentation",
    ], left=0.5, top=1.5, width=4, font_size=15)

    # Limitations
    add_bullet_content(slide, [
        ("Limitations", True, ACCENT_ORANGE, 22),
        "Single shared API key",
        "SQLite (not production-ready)",
        "Static seed data only",
        "No rate limiting",
        "No audio playback",
    ], left=4.8, top=1.5, width=3.5, font_size=15)

    # Future Work
    add_bullet_content(slide, [
        ("Future Work", True, MEDIUM_BLUE, 22),
        "OAuth2/JWT per-user auth",
        "PostgreSQL migration",
        "Spotify/MusicBrainz integration",
        "React frontend dashboard",
        "Docker containerisation",
        "Rate limiting & caching",
    ], left=8.8, top=1.5, width=4, font_size=15)

    # Thank you
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank you — Questions?"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = "docs/PRESENTATION.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
